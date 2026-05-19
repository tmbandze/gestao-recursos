from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Cores ──────────────────────────────────────────────
BG       = RGBColor(0x1E, 0x1E, 0x2E)   # fundo escuro
ACCENT   = RGBColor(0x4F, 0xC3, 0xF7)   # azul claro
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TXT = RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG  = RGBColor(0x2A, 0x2A, 0x3E)
YELLOW   = RGBColor(0xFF, 0xD7, 0x00)
GREEN    = RGBColor(0x98, 0xD8, 0x9E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout em branco

# ── Helpers ────────────────────────────────────────────

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_para(tf, text, size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def add_code_box(slide, lines, left, top, width, height, font_size=11.5):
    """Caixa com fundo escuro e texto monoespaçado."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x66)
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = False

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = GREEN

def add_title(slide, text, subtitle=None):
    add_textbox(slide, text,
                left=0.4, top=0.25, width=12.5, height=0.7,
                size=30, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    left=0.4, top=0.95, width=12.5, height=0.4,
                    size=16, color=GREY_TXT, italic=True)

def bullet(slide, items, left, top, width, height, size=16, gap=4):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(gap)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

# Linha decorativa topo
bar = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Título principal
tb = s1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Sistema de Gestão de Recursos"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE

# Subtítulo
tb2 = s1.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.7))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Partilha de Livros entre Estudantes"
r2.font.size = Pt(24); r2.font.color.rgb = ACCENT; r2.font.italic = True

# Linha separadora
sep = s1.shapes.add_shape(1, Inches(3.5), Inches(3.85), Inches(6.33), Inches(0.04))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0x44,0x44,0x66)
sep.line.fill.background()

# Disciplina
tb3 = s1.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.5))
tf3 = tb3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Sistemas Distribuídos — Trabalho Prático A"
r3.font.size = Pt(18); r3.font.color.rgb = GREY_TXT

tb4 = s1.shapes.add_textbox(Inches(1.5), Inches(4.65), Inches(10.3), Inches(0.4))
tf4 = tb4.text_frame
p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run(); r4.text = "Ano lectivo 2025/2026"
r4.font.size = Pt(16); r4.font.color.rgb = GREY_TXT

# Linha decorativa base
bar2 = s1.shapes.add_shape(1, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08))
bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT
bar2.line.fill.background()

# ══════════════════════════════════════════════════════
# SLIDE 2 — STACK TECNOLÓGICO
# ══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_title(s2, "Stack Tecnológico",
          "Duas JVMs independentes — servidor e cliente em processos separados")

# Tabela
rows, cols = 6, 3
table = s2.shapes.add_table(rows, cols,
    Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.6)).table

headers = ["Camada", "Tecnologia", "Versão"]
data = [
    ["Interface Gráfica",          "JavaFX",                        "23"],
    ["Linguagem",                  "Java SE",                       "17+"],
    ["Comunicação em rede",        "ServerSocket / Socket (TCP)",   "Java stdlib"],
    ["Serialização / Persistência","Gson  →  livros.json",          "2.10.1"],
    ["Build e dependências",       "Apache Maven (multi-módulo)",   "3.8+"],
]

col_widths = [Inches(3.2), Inches(5.8), Inches(2.7)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

def set_cell(cell, text, bg, fg, bold=False, size=16):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = fg
    cell.margin_top = Pt(6); cell.margin_bottom = Pt(6)

HDR_BG  = RGBColor(0x1A, 0x5C, 0x8A)
ROW_BG  = RGBColor(0x25, 0x25, 0x38)
ROW_ALT = RGBColor(0x2E, 0x2E, 0x45)

for c, h in enumerate(headers):
    set_cell(table.cell(0, c), h, HDR_BG, WHITE, bold=True, size=17)

for r, row in enumerate(data):
    bg = ROW_BG if r % 2 == 0 else ROW_ALT
    for c, val in enumerate(row):
        fg = ACCENT if c == 1 else WHITE
        set_cell(table.cell(r+1, c), val, bg, fg, size=16)

# ══════════════════════════════════════════════════════
# SLIDE 3 — COMUNICAÇÃO TCP
# ══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_title(s3, "Comunicação TCP: Sockets e Protocolo",
          '"Como cliente e servidor trocam mensagens?"')

# Diagrama de sequência (caixa de código)
diag = [
    "  Cliente                          Servidor",
    "    │                                 │",
    "    │── LOGIN|João ─────────────────► │",
    "    │◄─ OK|Bem-vindo ────────────────  │",
    "    │                                 │",
    "    │── REQUISITAR|uuid-1234 ────────► │",
    "    │◄─ OK|Livro requisitado ──────── │",
    "    │                                 │",
    "    │◄─ NOTIFICACAO|Disponível! ──────│   ← push assíncrono",
    "    │◄─ ATUALIZAR ───────────────────  │   ← broadcast a todos",
]
add_code_box(s3, diag, left=0.5, top=1.45, width=7.5, height=4.2, font_size=12)

# Bullets
bullet(s3, [
    "ServerSocket aguarda conexões na porta 8080",
    "Protocolo de texto simples:  COMANDO|campo1|campo2\\n",
    "Testável com telnet — independente da linguagem do cliente",
    "Funciona como middleware de abstracção",
], left=8.2, top=1.45, width=4.9, height=4.2, size=15)

# ══════════════════════════════════════════════════════
# SLIDE 4 — CONCORRÊNCIA
# ══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_title(s4, "Concorrência: ExecutorService + synchronized",
          '"Como o servidor atende vários clientes ao mesmo tempo?"')

code_pool = [
    "// Thread pool — uma thread por cliente",
    "ExecutorService executor = Executors.newCachedThreadPool();",
    "while (true) {",
    "    Socket socket = serverSocket.accept();",
    "    executor.submit(new GestorClientes(socket, gestorLivros));",
    "}",
]
add_code_box(s4, code_pool, left=0.4, top=1.45, width=6.3, height=2.3, font_size=12)

code_sync = [
    "// Secção crítica — synchronized",
    "public synchronized String requisitar(String id, ...) {",
    "    if (livro.isDisponivel()) {",
    "        livro.setEstudante(cliente.getNome());",
    '        return "OK|Livro requisitado";',
    "    } else {",
    "        livro.adicionarFila(cliente);",
    '        return "OK|Adicionado à fila";',
    "    }",
    "}",
]
add_code_box(s4, code_sync, left=0.4, top=3.9, width=6.3, height=3.0, font_size=12)

bullet(s4, [
    "newCachedThreadPool() — thread por cliente, sem bloquear os outros",
    "synchronized — secção crítica: só uma thread executa de cada vez",
    "Métodos de escrita sincronizados; leitura pura não precisa de lock",
    "Impossível dois clientes requisitarem o mesmo livro em simultâneo",
], left=6.9, top=1.45, width=6.2, height=5.5, size=15)

# ══════════════════════════════════════════════════════
# SLIDE 5 — NOTIFICAÇÕES ASSÍNCRONAS
# ══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_title(s5, "Notificações Assíncronas (Push)",
          '"O servidor avisa o cliente sem que ele tenha perguntado"')

# Caixa esquerda — polling (vermelho)
box_l = s5.shapes.add_shape(1, Inches(0.4), Inches(1.55), Inches(4.1), Inches(1.7))
box_l.fill.solid(); box_l.fill.fore_color.rgb = RGBColor(0x3E, 0x1E, 0x1E)
box_l.line.color.rgb = RGBColor(0xAA, 0x33, 0x33); box_l.line.width = Pt(1)
tb_l = box_l.text_frame; tb_l.word_wrap = True
p = tb_l.paragraphs[0]; r = p.add_run()
r.text = "Sem push — polling  ✗"; r.font.size = Pt(14); r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0x66, 0x66)
add_para(tb_l, '"Há novidades?" → "Há novidades?" → ...', 13, color=GREY_TXT, italic=True)
add_para(tb_l, "Cliente pergunta repetidamente — ineficiente", 13, color=RGBColor(0xFF,0x88,0x88))

# Caixa direita — push (verde)
box_r = s5.shapes.add_shape(1, Inches(4.7), Inches(1.55), Inches(4.1), Inches(1.7))
box_r.fill.solid(); box_r.fill.fore_color.rgb = RGBColor(0x1A, 0x33, 0x1A)
box_r.line.color.rgb = RGBColor(0x44, 0xAA, 0x44); box_r.line.width = Pt(1)
tb_r = box_r.text_frame; tb_r.word_wrap = True
p2 = tb_r.paragraphs[0]; r2 = p2.add_run()
r2.text = "Com push — este sistema  ✓"; r2.font.size = Pt(14); r2.font.bold = True
r2.font.color.rgb = GREEN
add_para(tb_r, "Servidor envia quando há novidade:", 13, color=GREY_TXT, italic=True)
add_para(tb_r, "NOTIFICACAO|Livro disponível para si!", 13, color=GREEN)

# Código NotificacaoService
code_notif = [
    "public void run() {",
    "    String msg;",
    "    while ((msg = entrada.readLine()) != null) {",
    '        if (msg.startsWith("NOTIFICACAO|"))',
    "            Platform.runLater(() -> controlador.mostrarNotificacao(msg));",
    '        else if (msg.equals("ATUALIZAR"))',
    "            Platform.runLater(() -> controlador.atualizarLista());",
    "        else",
    "            cliente.receberResposta(msg);  // resposta síncrona",
    "    }",
    "}",
]
add_code_box(s5, code_notif, left=0.4, top=3.4, width=8.6, height=3.6, font_size=11.5)

bullet(s5, [
    "Thread única a ler o socket — evita race condition",
    "Platform.runLater() — JavaFX exige actualizações na sua thread principal",
    "BlockingQueue sincroniza notificações com respostas síncronas",
], left=9.2, top=3.5, width=3.9, height=3.4, size=14)

# ══════════════════════════════════════════════════════
out = r"C:\dev\gestao-recursos\apresentacao_sd.pptx"
prs.save(out)
print(f"Ficheiro criado: {out}")
